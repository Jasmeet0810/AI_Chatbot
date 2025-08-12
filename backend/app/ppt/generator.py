from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
from typing import Dict, List, Optional, Any
from ..ai.content_generator import ContentGenerator
from ..config import settings
import logging

logger = logging.getLogger(__name__)

class PPTGenerator:
    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path or os.path.join(settings.template_dir, "lazulite_template.pptx")
        self.content_generator = ContentGenerator()
        self.output_dir = settings.generated_dir
        
        # Ensure output directory exists
        os.makedirs(self.output_dir, exist_ok=True)
    
    def generate_presentation(self, multi_product_data: List[Dict[str, Any]], user_prompt: str, user_id: str = None) -> str:
        """Generate PPT from multiple product data and user prompt"""
        try:
            logger.info(f"Starting multi-product PPT generation for {len(multi_product_data)} products")
            
            # Load template or create new presentation
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                logger.info(f"Loaded template: {self.template_path}")
            else:
                prs = Presentation()
                logger.info("Created new presentation (template not found)")
            
            # Generate AI-enhanced content for all products
            enhanced_content = self.content_generator.enhance_multi_product_content(multi_product_data, user_prompt)
            
            # Clear existing slides if using template
            if len(prs.slides) > 0:
                # Keep the first slide as title slide
                slides_to_remove = list(prs.slides)[1:]
                for slide in slides_to_remove:
                    rId = prs.slides._sldIdLst[prs.slides.index(slide)].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[prs.slides.index(slide)]
            
            # Create slides
            self.create_title_slide(prs, enhanced_content)
            
            # Create slides for each product
            for product_data in multi_product_data:
                self.create_product_overview_slide(prs, product_data)
                self.create_product_specifications_slide(prs, product_data)
                self.create_product_integration_slide(prs, product_data)
                self.create_product_infrastructure_slide(prs, product_data)
            
            self.create_conclusion_slide(prs, enhanced_content)
            
            # Generate unique filename
            safe_prompt = self._sanitize_filename(user_prompt)
            filename = f"presentation_{user_id or 'user'}_{safe_prompt}_{hash(user_prompt) % 10000}.pptx"
            output_path = os.path.join(self.output_dir, filename)
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"PPT saved successfully: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to generate presentation: {str(e)}")
            raise Exception(f"Failed to generate presentation: {str(e)}")
    
    def create_product_overview_slide(self, prs, product_data: Dict[str, Any]):
        """Create overview slide for a specific product"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = f"{product_data['product_name']} - Overview"
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = product_data.get('overview', 'Product overview not available')
        
        # Add images based on layout
        images = product_data.get('images', [])
        image_layout = product_data.get('image_layout', 'none')
        
        if images and image_layout != 'none':
            self.add_images_with_layout(slide, images, image_layout)
    
    def create_product_specifications_slide(self, prs, product_data: Dict[str, Any]):
        """Create specifications slide for a specific product"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = f"{product_data['product_name']} - Specifications"
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            specs_text = '\n'.join([f"• {spec}" for spec in product_data.get('specifications', [])])
            content_placeholder.text = specs_text or 'Specifications not available'
    
    def create_product_integration_slide(self, prs, product_data: Dict[str, Any]):
        """Create content integration slide for a specific product"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = f"{product_data['product_name']} - Content Integration"
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            integration_text = '\n'.join([f"• {item}" for item in product_data.get('content_integration', [])])
            content_placeholder.text = integration_text or 'Content integration information not available'
    
    def create_product_infrastructure_slide(self, prs, product_data: Dict[str, Any]):
        """Create infrastructure requirements slide for a specific product"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = f"{product_data['product_name']} - Infrastructure"
        
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            infra_text = '\n'.join([f"• {req}" for req in product_data.get('infrastructure_requirements', [])])
            content_placeholder.text = infra_text or 'Infrastructure requirements not available'
    
    def create_title_slide(self, prs, content):
        """Create title slide"""
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
        else:
            slide = prs.slides[0]
        
        # Update title and subtitle
        if slide.shapes.title:
            slide.shapes.title.text = content.get('title', 'Lazulite Product Presentation')
        
        # Find subtitle placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = content.get('subtitle', 'Multi-Product Presentation - Powered by Lazulite AI')
                break
    
    def create_conclusion_slide(self, prs, content):
        """Create conclusion slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Thank You"
        
        conclusion_text = "Thank you for your interest in Lazulite products!\n\nFor more information, visit: https://lazulite.ae"
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = conclusion_text
    
    def add_images_with_layout(self, slide, images: List[str], layout: str):
        """Add images to slide based on specified layout"""
        try:
            if not images:
                return
            
            if layout == "single":
                self.add_single_image(slide, images[0])
            elif layout == "side_by_side":
                self.add_side_by_side_images(slide, images[:2])
            elif layout == "grid":
                self.add_grid_images(slide, images[:4])  # Max 4 for grid
                
        except Exception as e:
            logger.error(f"Failed to add images with layout {layout}: {str(e)}")
    
    def add_single_image(self, slide, image_path: str):
        """Add single centered image"""
        try:
            if os.path.exists(image_path):
                left = Inches(2)
                top = Inches(3)
                width = Inches(6)
                height = Inches(4)
                slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            logger.warning(f"Failed to add single image: {str(e)}")
    
    def add_side_by_side_images(self, slide, images: List[str]):
        """Add two images side by side"""
        try:
            img_width = Inches(3)
            img_height = Inches(2.5)
            
            for i, img_path in enumerate(images[:2]):
                if os.path.exists(img_path):
                    left = Inches(1 + i * 4)
                    top = Inches(3.5)
                    slide.shapes.add_picture(img_path, left, top, img_width, img_height)
        except Exception as e:
            logger.warning(f"Failed to add side by side images: {str(e)}")
    
    def add_grid_images(self, slide, images: List[str]):
        """Add images in 2x2 grid"""
        try:
            img_width = Inches(2.5)
            img_height = Inches(1.8)
            
            for i, img_path in enumerate(images[:4]):
                if os.path.exists(img_path):
                    col = i % 2
                    row = i // 2
                    left = Inches(1 + col * 3.5)
                    top = Inches(3 + row * 2.2)
                    slide.shapes.add_picture(img_path, left, top, img_width, img_height)
        except Exception as e:
            logger.warning(f"Failed to add grid images: {str(e)}")
    
    def add_images_to_slide(self, slide, images: List[str]):
        """Add converted images to slide"""
        try:
            # Position images in a grid
            img_width = Inches(2.5)
            img_height = Inches(1.8)
            
            for i, img_path in enumerate(images):
                if os.path.exists(img_path):
                    # Calculate position (2x2 grid)
                    col = i % 2
                    row = i // 2
                    left = Inches(0.5 + col * 3)
                    top = Inches(3 + row * 2.2)
                    
                    try:
                        slide.shapes.add_picture(img_path, left, top, img_width, img_height)
                        logger.info(f"Added image to slide: {img_path}")
                    except Exception as e:
                        logger.warning(f"Failed to add image {img_path}: {str(e)}")
                        
        except Exception as e:
            logger.error(f"Failed to add images to slide: {str(e)}")
    
    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for safe file operations"""
        import re
        if not filename:
            return "presentation"
        
        # Remove or replace invalid characters
        filename = re.sub(r'[^\w\s-]', '', filename)
        filename = re.sub(r'[-\s]+', '_', filename)
        
        return filename.lower()[:30]  # Limit length
    
    def get_presentation_info(self, file_path: str) -> Dict:
        """Get information about generated presentation"""
        try:
            if not os.path.exists(file_path):
                return {"error": "File not found"}
            
            prs = Presentation(file_path)
            return {
                "slide_count": len(prs.slides),
                "file_size": os.path.getsize(file_path),
                "file_path": file_path,
                "filename": os.path.basename(file_path)
            }
        except Exception as e:
            logger.error(f"Failed to get presentation info: {str(e)}")
            return {"error": str(e)}