from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)

class LazuliteStyleGuide:
    """Lazulite brand style guide for presentations"""
    
    # Brand Colors
    BRAND_COLORS = {
        # Primary Colors
        'lazulite_blue': RGBColor(30, 64, 175),      # #1E40AF
        'lazulite_purple': RGBColor(124, 58, 237),   # #7C3AED  
        'lazulite_teal': RGBColor(13, 148, 136),     # #0D9488
        
        # Secondary Colors
        'light_blue': RGBColor(59, 130, 246),        # #3B82F6
        'light_purple': RGBColor(147, 51, 234),      # #9333EA
        'light_teal': RGBColor(20, 184, 166),        # #14B8A6
        
        # Neutral Colors
        'dark_gray': RGBColor(31, 41, 55),           # #1F2937
        'medium_gray': RGBColor(107, 114, 128),      # #6B7280
        'light_gray': RGBColor(156, 163, 175),       # #9CA3AF
        'very_light_gray': RGBColor(243, 244, 246),  # #F3F4F6
        'white': RGBColor(255, 255, 255),            # #FFFFFF
        'black': RGBColor(0, 0, 0),                  # #000000
    }
    
    # Typography
    TYPOGRAPHY = {
        'primary_font': 'Inter',
        'secondary_font': 'Calibri',
        'fallback_font': 'Arial',
        
        # Font Sizes
        'title_size': Pt(48),
        'subtitle_size': Pt(28),
        'heading_1_size': Pt(36),
        'heading_2_size': Pt(28),
        'heading_3_size': Pt(24),
        'body_size': Pt(20),
        'caption_size': Pt(16),
        'small_size': Pt(14),
    }
    
    # Spacing and Layout
    LAYOUT = {
        'slide_margin': 0.5,  # inches
        'content_margin': 1.0,  # inches
        'line_spacing': 1.2,
        'paragraph_spacing': Pt(12),
        'bullet_indent': 0.25,  # inches
    }

class PPTStyler:
    """Apply consistent styling to PowerPoint presentations"""
    
    def __init__(self):
        self.style_guide = LazuliteStyleGuide()
    
    def apply_slide_styles(self, slide, slide_type: str = 'content'):
        """Apply consistent styling to a slide"""
        try:
            if slide_type == 'title':
                self._style_title_slide(slide)
            elif slide_type == 'content':
                self._style_content_slide(slide)
            elif slide_type == 'section':
                self._style_section_slide(slide)
            else:
                self._style_content_slide(slide)  # Default
                
        except Exception as e:
            logger.error(f"Failed to apply slide styles: {str(e)}")
    
    def _style_title_slide(self, slide):
        """Style title slide"""
        try:
            # Style title
            if slide.shapes.title:
                title_frame = slide.shapes.title.text_frame
                self._apply_title_style(title_frame)
            
            # Style subtitle
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Subtitle
                    self._apply_subtitle_style(shape.text_frame)
                    break
                    
        except Exception as e:
            logger.warning(f"Failed to style title slide: {str(e)}")
    
    def _style_content_slide(self, slide):
        """Style content slide"""
        try:
            # Style title
            if slide.shapes.title:
                title_frame = slide.shapes.title.text_frame
                self._apply_heading_style(title_frame, level=1)
            
            # Style content
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content
                    self._apply_body_style(shape.text_frame)
                    break
                    
        except Exception as e:
            logger.warning(f"Failed to style content slide: {str(e)}")
    
    def _style_section_slide(self, slide):
        """Style section divider slide"""
        try:
            # Style as title slide but with different colors
            if slide.shapes.title:
                title_frame = slide.shapes.title.text_frame
                self._apply_section_title_style(title_frame)
                
        except Exception as e:
            logger.warning(f"Failed to style section slide: {str(e)}")
    
    def _apply_title_style(self, text_frame):
        """Apply title styling"""
        try:
            # Clear existing paragraphs except first
            while len(text_frame.paragraphs) > 1:
                text_frame._element.remove(text_frame.paragraphs[-1]._element)
            
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_after = self.style_guide.LAYOUT['paragraph_spacing']
            
            # Apply font styling
            if paragraph.runs:
                run = paragraph.runs[0]
            else:
                run = paragraph.add_run()
            
            run.font.name = self.style_guide.TYPOGRAPHY['primary_font']
            run.font.size = self.style_guide.TYPOGRAPHY['title_size']
            run.font.color.rgb = self.style_guide.BRAND_COLORS['lazulite_blue']
            run.font.bold = True
            
        except Exception as e:
            logger.warning(f"Failed to apply title style: {str(e)}")
    
    def _apply_subtitle_style(self, text_frame):
        """Apply subtitle styling"""
        try:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            if paragraph.runs:
                run = paragraph.runs[0]
            else:
                run = paragraph.add_run()
            
            run.font.name = self.style_guide.TYPOGRAPHY['secondary_font']
            run.font.size = self.style_guide.TYPOGRAPHY['subtitle_size']
            run.font.color.rgb = self.style_guide.BRAND_COLORS['medium_gray']
            run.font.italic = True
            
        except Exception as e:
            logger.warning(f"Failed to apply subtitle style: {str(e)}")
    
    def _apply_heading_style(self, text_frame, level: int = 1):
        """Apply heading styling"""
        try:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            
            if paragraph.runs:
                run = paragraph.runs[0]
            else:
                run = paragraph.add_run()
            
            # Choose size and color based on level
            if level == 1:
                size = self.style_guide.TYPOGRAPHY['heading_1_size']
                color = self.style_guide.BRAND_COLORS['lazulite_purple']
            elif level == 2:
                size = self.style_guide.TYPOGRAPHY['heading_2_size']
                color = self.style_guide.BRAND_COLORS['lazulite_teal']
            else:
                size = self.style_guide.TYPOGRAPHY['heading_3_size']
                color = self.style_guide.BRAND_COLORS['dark_gray']
            
            run.font.name = self.style_guide.TYPOGRAPHY['primary_font']
            run.font.size = size
            run.font.color.rgb = color
            run.font.bold = True
            
        except Exception as e:
            logger.warning(f"Failed to apply heading style: {str(e)}")
    
    def _apply_body_style(self, text_frame):
        """Apply body text styling"""
        try:
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_after = Pt(6)
                paragraph.line_spacing = self.style_guide.LAYOUT['line_spacing']
                
                for run in paragraph.runs:
                    run.font.name = self.style_guide.TYPOGRAPHY['secondary_font']
                    run.font.size = self.style_guide.TYPOGRAPHY['body_size']
                    run.font.color.rgb = self.style_guide.BRAND_COLORS['dark_gray']
                    
        except Exception as e:
            logger.warning(f"Failed to apply body style: {str(e)}")
    
    def _apply_section_title_style(self, text_frame):
        """Apply section title styling"""
        try:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            if paragraph.runs:
                run = paragraph.runs[0]
            else:
                run = paragraph.add_run()
            
            run.font.name = self.style_guide.TYPOGRAPHY['primary_font']
            run.font.size = self.style_guide.TYPOGRAPHY['heading_1_size']
            run.font.color.rgb = self.style_guide.BRAND_COLORS['lazulite_teal']
            run.font.bold = True
            
        except Exception as e:
            logger.warning(f"Failed to apply section title style: {str(e)}")
    
    def create_bullet_points(self, text_frame, items: list):
        """Create styled bullet points"""
        try:
            # Clear existing content
            text_frame.clear()
            
            for i, item in enumerate(items):
                if i == 0:
                    paragraph = text_frame.paragraphs[0]
                else:
                    paragraph = text_frame.add_paragraph()
                
                paragraph.text = str(item)
                paragraph.level = 0
                paragraph.space_after = Pt(6)
                
                # Style the run
                if paragraph.runs:
                    run = paragraph.runs[0]
                    run.font.name = self.style_guide.TYPOGRAPHY['secondary_font']
                    run.font.size = self.style_guide.TYPOGRAPHY['body_size']
                    run.font.color.rgb = self.style_guide.BRAND_COLORS['dark_gray']
                    
        except Exception as e:
            logger.warning(f"Failed to create bullet points: {str(e)}")
    
    def apply_gradient_background(self, slide, gradient_type: str = 'blue_purple'):
        """Apply gradient background to slide"""
        try:
            # Note: Gradient backgrounds are complex in python-pptx
            # This is a simplified implementation
            
            # Add a rectangle shape covering the entire slide
            from pptx.util import Inches
            
            # Get slide dimensions (standard 16:9)
            slide_width = Inches(13.33)
            slide_height = Inches(7.5)
            
            # Add background shape
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0, slide_width, slide_height
            )
            
            # Move to back
            slide.shapes._spTree.remove(background._element)
            slide.shapes._spTree.insert(2, background._element)
            
            # Apply fill color (simplified - no actual gradient)
            if gradient_type == 'blue_purple':
                fill_color = self.style_guide.BRAND_COLORS['lazulite_blue']
            elif gradient_type == 'teal_blue':
                fill_color = self.style_guide.BRAND_COLORS['lazulite_teal']
            else:
                fill_color = self.style_guide.BRAND_COLORS['very_light_gray']
            
            background.fill.solid()
            background.fill.fore_color.rgb = fill_color
            
        except Exception as e:
            logger.warning(f"Failed to apply gradient background: {str(e)}")
    
    def get_style_config(self) -> Dict[str, Any]:
        """Get complete style configuration"""
        return {
            'colors': {name: f"#{color.rgb:06x}" for name, color in self.style_guide.BRAND_COLORS.items()},
            'fonts': self.style_guide.TYPOGRAPHY,
            'layout': self.style_guide.LAYOUT
        }