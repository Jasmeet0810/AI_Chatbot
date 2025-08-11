import os
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from ..config import settings
import logging

logger = logging.getLogger(__name__)

class PPTTemplateManager:
    def __init__(self):
        self.template_dir = settings.template_dir
        self.default_template = "lazulite_template.pptx"
        
        # Ensure template directory exists
        os.makedirs(self.template_dir, exist_ok=True)
    
    def get_template_path(self, template_name: str = None) -> str:
        """Get path to template file"""
        template_name = template_name or self.default_template
        return os.path.join(self.template_dir, template_name)
    
    def create_default_template(self) -> str:
        """Create default Lazulite template if it doesn't exist"""
        template_path = self.get_template_path()
        
        if os.path.exists(template_path):
            logger.info(f"Template already exists: {template_path}")
            return template_path
        
        try:
            logger.info("Creating default Lazulite template")
            
            # Create new presentation
            prs = Presentation()
            
            # Customize slide layouts
            self._customize_title_layout(prs)
            self._customize_content_layout(prs)
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Set default title and subtitle
            slide.shapes.title.text = "Lazulite Product Presentation"
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Powered by Lazulite AI Technology"
            
            # Save template
            prs.save(template_path)
            logger.info(f"Default template created: {template_path}")
            
            return template_path
            
        except Exception as e:
            logger.error(f"Failed to create default template: {str(e)}")
            raise Exception(f"Template creation failed: {str(e)}")
    
    def _customize_title_layout(self, prs):
        """Customize title slide layout with Lazulite branding"""
        try:
            title_layout = prs.slide_layouts[0]
            
            # Customize background (if possible)
            # Note: Background customization is limited in python-pptx
            
            # Set Lazulite color scheme
            # Blue: #1E40AF, Purple: #7C3AED, Teal: #0D9488
            
        except Exception as e:
            logger.warning(f"Failed to customize title layout: {str(e)}")
    
    def _customize_content_layout(self, prs):
        """Customize content slide layout"""
        try:
            content_layout = prs.slide_layouts[1]
            
            # Add Lazulite branding elements if needed
            
        except Exception as e:
            logger.warning(f"Failed to customize content layout: {str(e)}")
    
    def validate_template(self, template_path: str) -> Dict:
        """Validate template file"""
        validation_result = {
            "is_valid": True,
            "errors": [],
            "warnings": [],
            "info": {}
        }
        
        try:
            if not os.path.exists(template_path):
                validation_result["is_valid"] = False
                validation_result["errors"].append("Template file not found")
                return validation_result
            
            # Try to open template
            prs = Presentation(template_path)
            
            # Check basic requirements
            if len(prs.slide_layouts) < 2:
                validation_result["warnings"].append("Template has fewer than 2 slide layouts")
            
            validation_result["info"] = {
                "slide_layouts": len(prs.slide_layouts),
                "file_size": os.path.getsize(template_path),
                "filename": os.path.basename(template_path)
            }
            
            logger.info(f"Template validation successful: {template_path}")
            
        except Exception as e:
            validation_result["is_valid"] = False
            validation_result["errors"].append(f"Template validation failed: {str(e)}")
            logger.error(f"Template validation failed: {str(e)}")
        
        return validation_result
    
    def list_available_templates(self) -> List[Dict]:
        """List all available templates"""
        templates = []
        
        try:
            if not os.path.exists(self.template_dir):
                return templates
            
            for filename in os.listdir(self.template_dir):
                if filename.endswith('.pptx'):
                    template_path = os.path.join(self.template_dir, filename)
                    validation = self.validate_template(template_path)
                    
                    template_info = {
                        "name": filename,
                        "path": template_path,
                        "is_valid": validation["is_valid"],
                        "is_default": filename == self.default_template
                    }
                    
                    if validation["info"]:
                        template_info.update(validation["info"])
                    
                    templates.append(template_info)
            
        except Exception as e:
            logger.error(f"Failed to list templates: {str(e)}")
        
        return templates
    
    def ensure_template_exists(self) -> str:
        """Ensure default template exists, create if necessary"""
        template_path = self.get_template_path()
        
        if not os.path.exists(template_path):
            return self.create_default_template()
        
        # Validate existing template
        validation = self.validate_template(template_path)
        if not validation["is_valid"]:
            logger.warning("Existing template is invalid, creating new one")
            # Backup old template
            backup_path = template_path + ".backup"
            if os.path.exists(template_path):
                os.rename(template_path, backup_path)
            
            return self.create_default_template()
        
        return template_path

class PPTStyleManager:
    """Manage PPT styling and formatting"""
    
    # Lazulite color palette
    COLORS = {
        'primary_blue': RGBColor(30, 64, 175),      # #1E40AF
        'primary_purple': RGBColor(124, 58, 237),   # #7C3AED
        'primary_teal': RGBColor(13, 148, 136),     # #0D9488
        'text_dark': RGBColor(31, 41, 55),          # #1F2937
        'text_light': RGBColor(107, 114, 128),      # #6B7280
        'background': RGBColor(249, 250, 251),      # #F9FAFB
        'white': RGBColor(255, 255, 255),           # #FFFFFF
    }
    
    # Font settings
    FONTS = {
        'title': 'Calibri',
        'heading': 'Calibri',
        'body': 'Calibri',
        'caption': 'Calibri'
    }
    
    # Font sizes
    FONT_SIZES = {
        'title': Pt(44),
        'subtitle': Pt(24),
        'heading': Pt(32),
        'subheading': Pt(24),
        'body': Pt(18),
        'caption': Pt(14)
    }
    
    @classmethod
    def apply_title_style(cls, text_frame):
        """Apply title styling"""
        try:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = cls.FONTS['title']
            run.font.size = cls.FONT_SIZES['title']
            run.font.color.rgb = cls.COLORS['primary_blue']
            run.font.bold = True
            
        except Exception as e:
            logger.warning(f"Failed to apply title style: {str(e)}")
    
    @classmethod
    def apply_heading_style(cls, text_frame):
        """Apply heading styling"""
        try:
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = cls.FONTS['heading']
            run.font.size = cls.FONT_SIZES['heading']
            run.font.color.rgb = cls.COLORS['primary_purple']
            run.font.bold = True
            
        except Exception as e:
            logger.warning(f"Failed to apply heading style: {str(e)}")
    
    @classmethod
    def apply_body_style(cls, text_frame):
        """Apply body text styling"""
        try:
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                
                for run in paragraph.runs:
                    run.font.name = cls.FONTS['body']
                    run.font.size = cls.FONT_SIZES['body']
                    run.font.color.rgb = cls.COLORS['text_dark']
                    
        except Exception as e:
            logger.warning(f"Failed to apply body style: {str(e)}")